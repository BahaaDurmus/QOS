"""
QOS_Sunum.pptx olusturur. Duzenlemek icin PowerPoint/LibreOffice ile acin.
  python presentation/generate_pptx.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUT = Path(__file__).parent / "QOS_Sunum.pptx"

SLIDES = [
    ("lead", "Kablosuz Örgüsel Ağlarda\nQoS Tabanlı Video Akışı Optimizasyonu",
     "Gerçek zamanlı UDP · WMN simülasyonu · VSR onarımı\n\nAdınız Soyadınız · Bölüm · Tarih"),
    ("bullet", "Problem",
     ["WMN/WiFi: paket kaybı, gecikme, jitter",
      "Eksik UDP → JPEG karede siyah bantlar",
      "Sadece bitrate düşürmek görüntüyü bulanıklaştırır",
      "Hedef: QoS izleme + yapay zeka ile onarım"]),
    ("bullet", "Proje hedefi",
     ["WMN benzeri ağ simülasyonu (proxy)",
      "QoS metrikleri → ABR kalite kademesi",
      "15 kare VSR modeli ile bozuk bölgeleri onarma",
      "Web dashboard ile canlı metrik ve görüntü izleme"]),
    ("bullet", "Mimari — 4 aşama",
     ["1. stream_server.py — MP4 → JPEG → UDP :9998",
      "2. wmn_simulator.py — kayıp/gecikme → :9999",
      "3. qos_monitor.py — HIGH/MEDIUM/LOW/CRITICAL",
      "4. stream_client.py — VSR + canlı önizleme"]),
    ("bullet", "Aşama 1 — Video sunucusu",
     ["JPEG kalite + FPS ayarlı gönderim",
      "60 KB chunk, 12 bayt header",
      "frame_id | total_chunks | chunk_id | data_len"]),
    ("table", "WMN profilleri",
     [["Profil", "Kayıp", "Gecikme"],
      ["medium", "~%5", "50 ms"],
      ["poor", "~%15", "150 ms"],
      ["critical", "~%30", "300 ms"]]),
    ("bullet", "Aşama 3–4 — QoS ve VSR",
     ["LOW/CRITICAL → VSR aktif",
      "VSRModel: 15×3 kanal giriş, 8 residual block",
      "Async thread — akış bloklanmaz",
      "model_sharp_ep10_fixed.pth (~2.4 MB)"]),
    ("bullet", "Web dashboard",
     ["http://localhost:8080 — canlı metrik + 2 video",
      "/pipeline.html — paket sayaçları, log",
      "/comparison.html — bozuk vs onarılmış (frame)"]),
    ("code", "Canlı demo",
     "cd QOS-master\npip install torch torchvision opencv-python pillow numpy\npython run_full_demo.py --video input.mp4 --model model_sharp_ep10_fixed.pth --profile medium"),
    ("table", "Ölçüm sonuçları (örnek)",
     [["Metrik", "Bozuk", "VSR", "Kazanım"],
      ["PSNR", "22.02 dB", "25.81 dB", "+3.79 dB"],
      ["SSIM", "0.7706", "0.9137", "+0.143"]]),
    ("bullet", "Güçlü yönler / sınırlamalar",
     ["✓ Uçtan uca prototip, modüler kod",
      "✓ Ölçülebilir PSNR/SSIM",
      "△ Yerel loopback (gerçek mesh değil)",
      "△ VSR gecikmesi (15 kare penceresi)"]),
    ("bullet", "Gelecek çalışmalar (düzenleyin)",
     ["Gerçek WMN test yatağı",
      "Hafif model / GPU hızlandırma",
      "QoS geri bildirimi ile bitrate kontrolü",
      "Çoklu istemci senaryoları"]),
    ("bullet", "Sorun giderme",
     ["Port meşgul → taskkill /F /IM python.exe",
      "Görüntü yok → 15–20 sn VSR ısınması",
      "comparison boş → make_web_frames.py"]),
    ("lead", "Teşekkürler", "Sorular?\n\ne-posta@universite.edu"),
]


def add_title_slide(prs, title, subtitle):
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if slide.placeholders[1]:
        slide.placeholders[1].text = subtitle


def add_bullet_slide(prs, title, bullets):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, line in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(20)


def add_code_slide(prs, title, code):
    layout = prs.slide_layouts[5]  # blank
    slide = prs.slides.add_slide(layout)
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tx.text_frame.text = title
    tx.text_frame.paragraphs[0].font.size = Pt(32)
    tx.text_frame.paragraphs[0].font.bold = True
    box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
    tf = box.text_frame
    tf.text = code
    for p in tf.paragraphs:
        p.font.name = "Consolas"
        p.font.size = Pt(14)


def add_table_slide(prs, title, rows):
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tx.text_frame.text = title
    tx.text_frame.paragraphs[0].font.size = Pt(32)
    tx.text_frame.paragraphs[0].font.bold = True
    nrows, ncols = len(rows), len(rows[0])
    tbl = slide.shapes.add_table(nrows, ncols, Inches(0.8), Inches(1.4), Inches(8.4), Inches(0.4 * nrows)).table
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(18 if r == 0 else 16)
                if r == 0:
                    p.font.bold = True


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for kind, title, content in SLIDES:
        if kind == "lead":
            add_title_slide(prs, title.replace("\n", " "), content)
        elif kind == "bullet":
            add_bullet_slide(prs, title, content)
        elif kind == "code":
            add_code_slide(prs, title, content)
        elif kind == "table":
            add_table_slide(prs, title, content)

    prs.save(OUT)
    print(f"[OK] Kaydedildi: {OUT.resolve()}")


if __name__ == "__main__":
    main()
